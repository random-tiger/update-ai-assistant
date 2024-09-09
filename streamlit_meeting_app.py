import streamlit as st
import os
from openai import OpenAI
from docx import Document
from io import BytesIO
import moviepy.editor as mp
import tempfile
import docx
import pandas as pd
import fitz
import base64
import requests
from st_aggrid import AgGrid, GridOptionsBuilder, GridUpdateMode
from streamlit_quill import st_quill
from pptx import Presentation
from PIL import Image
from concurrent.futures import ThreadPoolExecutor
import openpyxl
from openpyxl.drawing.image import Image as OpenpyxlImage
from pre_canned_prompts_file import pre_canned_prompts
from openpyxl import load_workbook
from defusedxml.common import DefusedXmlException
import defusedxml.ElementTree as ET

api_key = st.secrets["OPENAI_API_KEY"]

# Initialize OpenAI client
client = OpenAI(api_key=api_key)

def transcribe_audio(audio_file):
    """Transcribe audio using Whisper model."""
    transcription = client.audio.transcriptions.create(model="whisper-1", file=audio_file)
    return transcription['text'] if isinstance(transcription, dict) else transcription.text

def generate_response(transcription, model, custom_prompt):
    """Generate AI response based on the provided transcription and model."""
    response = client.chat.completions.create(
        model=model,
        temperature=0,
        messages=[
            {"role": "system", "content": custom_prompt},
            {"role": "user", "content": transcription}
        ]
    )
    return response.choices[0].message.content

def save_as_docx(minutes):
    """Save the generated meeting minutes as a Word document."""
    doc = Document()
    for key, value in minutes.items():
        doc.add_heading(key.replace('_', ' ').title(), level=1)
        doc.add_paragraph(value)
    buffer = BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer

def convert_video_to_mp3(uploaded_file, suffix):
    """Convert video files to MP3 format."""
    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as temp_video_file:
        temp_video_file.write(uploaded_file.getbuffer())
        temp_video_file_path = temp_video_file.name

    video = mp.VideoFileClip(temp_video_file_path)

    if not video.audio:
        st.error(f"The uploaded {suffix} file does not contain an audio track.")
        return None

    with tempfile.NamedTemporaryFile(delete=False, suffix=".mp3") as audio_file:
        video.audio.write_audiofile(audio_file.name)
        return audio_file.name

def encode_image(image):
    """Encode an image to Base64 format."""
    with BytesIO() as buffer:
        image.save(buffer, format=image.format)
        return base64.b64encode(buffer.getvalue()).decode()

def transcribe_image(image_file):
    """Transcribe image using GPT-4's multimodal capabilities."""
    base64_image = encode_image(image_file)

    headers = {
        "Content-Type": "application/json",
        "Authorization": f"Bearer {os.getenv('OPENAI_API_KEY')}"
    }

    payload = {
        "model": "gpt-4o-mini",
        "messages": [
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": "You describe images to vision-impaired people to help them understand the specific, detailed contents and meaning of them. Please translate and describe the details and likely meaning of this image."},
                    {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{base64_image}"}}
                ]
            }
        ],
        "max_tokens": 300
    }

    response = requests.post("https://api.openai.com/v1/chat/completions", headers=headers, json=payload)
    return response.json()['choices'][0]['message']['content']

def process_images_concurrently(images):
    """Process and transcribe images concurrently."""
    with ThreadPoolExecutor(max_workers=10) as executor:
        transcriptions = list(executor.map(transcribe_image, images))
    return transcriptions

def read_docx(file):
    """Read text and images from a DOCX file."""
    doc = docx.Document(file)
    result = ""
    images = []

    for para in doc.paragraphs:
        result += para.text + "\n"

    # Loop through all elements in the document to find images
    for rel in doc.part.rels.values():
        if "image" in rel.target_ref:
            image_data = rel.target_part.blob  # Retrieve image binary data
            image = Image.open(BytesIO(image_data))
            images.append(image)

    if images:
        image_transcriptions = process_images_concurrently(images)
        for transcription in image_transcriptions:
            result += f"\n[Image: {transcription}]\n"

    return result

def read_pdf(file):
    """Read text and images from a PDF file."""
    document = fitz.open(stream=file.read(), filetype="pdf")
    result = ""
    images = []

    for page_num in range(len(document)):
        page = document.load_page(page_num)
        page_text = page.get_text()
        result += f"\nPage {page_num + 1}\n{page_text}"

        page_images = []
        for img in page.get_images(full=True):
            xref = img[0]
            base_image = document.extract_image(xref)
            image = Image.open(BytesIO(base_image["image"]))
            page_images.append(image)

        if page_images:
            image_transcriptions = process_images_concurrently(page_images)
            for transcription in image_transcriptions:
                result += f"\n[Image on page {page_num + 1}: {transcription}]\n"

    return result

def read_pptx(file):
    """Read text and images from a PowerPoint file."""
    presentation = Presentation(file)
    result = ""

    for slide_num, slide in enumerate(presentation.slides, start=1):
        slide_text = ""
        images = []

        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text += shape.text + "\n"

            # Ensure to correctly handle image shapes
            if hasattr(shape, "image"):
                image = shape.image
                image_stream = image.blob  # Retrieve image binary data
                image = Image.open(BytesIO(image_stream))
                images.append(image)

        result += f"Slide {slide_num}:\n{slide_text}"

        if images:
            image_transcriptions = process_images_concurrently(images)
            for transcription in image_transcriptions:
                result += f"\n[Image: {transcription}]\n"

    return result

def read_txt(file):
    """Read text from a TXT file."""
    return file.read().decode("utf-8")

def read_excel(file):
    """Read text and images from an Excel file."""
    wb = openpyxl.load_workbook(file)
    result = ""
    images = []

    for sheet in wb.sheetnames:
        ws = wb[sheet]
        result += f"Sheet: {sheet}\n"

        # Reading text content
        for row in ws.iter_rows(values_only=True):
            result += "\t".join([str(cell) if cell is not None else "" for cell in row]) + "\n"

        # Reading images
        for img in ws._images:
            img_stream = img._data()  # Retrieve image binary data
            image = Image.open(BytesIO(img_stream))
            images.append(image)

    if images:
        image_transcriptions = process_images_concurrently(images)
        for transcription in image_transcriptions:
            result += f"\n[Image: {transcription}]\n"

    return result


def main():
    st.markdown(
        """
        <style>
        .reportview-container .main .block-container {
            padding-left: 0rem;
            padding-right: 0rem;
            max-width: 100%;
            margin: 0 auto.
        }
        .css-18e3th9 {
            flex: 1 1 100%;
            width: 100%;
            padding: 2rem 1rem 1rem;
        }
        </style>
        """, unsafe_allow_html=True
    )

    st.sidebar.title("Wonk")
    st.sidebar.info("Upload mp3, mp4, mov, docx, txt, xlsx, pdf, pptx, or image files to start!")
    uploaded_files = st.sidebar.file_uploader("Upload audio, video, text, or image files", type=["mp3", "mp4", "mov", "docx", "txt", "xlsx", "pdf", "pptx", "jpg", "jpeg", "png"], accept_multiple_files=True)
    process_files = st.sidebar.button("Process Files")

    if uploaded_files and process_files:
        st.session_state.setdefault('transcriptions', [])
        progress_bar = st.progress(0)
        total_files = len(uploaded_files)

        status_placeholder = st.empty()

        # Stage 1: Files Uploaded
        status_placeholder.info("Files uploaded successfully!")
        progress_bar.progress(0.1)

        # Initial processing stage
        for i, uploaded_file in enumerate(uploaded_files):
            with status_placeholder.container():
                status_placeholder.info(f"Processing {uploaded_file.name} ({i + 1}/{total_files})...")

            with st.spinner(f"Processing {uploaded_file.name}..."):
                # Handle different file types and process accordingly
                if uploaded_file.type in ["video/quicktime", "video/mp4"]:
                    status_placeholder.info("Converting video to audio...")
                    progress_bar.progress(0.2)
                    audio_file_path = convert_video_to_mp3(uploaded_file, ".mov" if uploaded_file.type == "video/quicktime" else ".mp4")
                    if audio_file_path:
                        with open(audio_file_path, "rb") as f:
                            status_placeholder.info("Transcribing audio...")
                            progress_bar.progress(0.5)
                            st.session_state.transcriptions.append(transcribe_audio(f))
                elif uploaded_file.type == "audio/mpeg":
                    status_placeholder.info("Transcribing audio file...")
                    progress_bar.progress(0.5)
                    st.session_state.transcriptions.append(transcribe_audio(uploaded_file))
                elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                    status_placeholder.info("Reading DOCX file...")
                    progress_bar.progress(0.2)
                    st.session_state.transcriptions.append(read_docx(uploaded_file))
                elif uploaded_file.type == "text/plain":
                    status_placeholder.info("Reading TXT file...")
                    progress_bar.progress(0.2)
                    st.session_state.transcriptions.append(read_txt(uploaded_file))
                elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
                    status_placeholder.info("Reading Excel file...")
                    progress_bar.progress(0.2)
                    st.session_state.transcriptions.append(read_excel(uploaded_file))
                elif uploaded_file.type == "application/pdf":
                    status_placeholder.info("Reading PDF file...")
                    progress_bar.progress(0.2)
                    st.session_state.transcriptions.append(read_pdf(uploaded_file))
                elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                    status_placeholder.info("Reading PowerPoint file...")
                    progress_bar.progress(0.2)
                    st.session_state.transcriptions.append(read_pptx(uploaded_file))
                elif uploaded_file.type in ["image/jpeg", "image/png"]:
                    status_placeholder.info("Transcribing image...")
                    progress_bar.progress(0.5)
                    st.session_state.transcriptions.append(transcribe_image(Image.open(uploaded_file)))

            # Update the progress bar after processing each file
            progress_bar.progress(0.2 + (0.7 * (i + 1) / total_files))

        # Stage 3: Processing Complete
        status_placeholder.success("All files processed successfully!")
        progress_bar.progress(1.0)

        if st.session_state.transcriptions:
            st.session_state.transcription = "\n\n".join(st.session_state.transcriptions)

    if "transcription" in st.session_state:
        with st.expander("Transcription", expanded=True):
            st.subheader("Transcription")
            st.session_state.transcription = st_quill(value=st.session_state.transcription, key='transcription_editor')

        st.sidebar.info("Select what you'd like to create!")
        summary_type = st.sidebar.radio("Select the type of summary you want to generate:", ["Meeting Summary", "User Research", "Action Items", "Retro", "Document Review"], index=0)

        st.session_state.setdefault('prompts', [])

        checkboxes = {}
        if summary_type:
            st.sidebar.markdown(f"### {summary_type} Prompts")
            st.sidebar.info("Select the sections you'd like in your document!")
            checkboxes = {key: st.sidebar.checkbox(heading["heading"]) for key, heading in pre_canned_prompts[summary_type.lower().replace(" ", "_")].items()}

        if any(checkboxes.values()):
            st.sidebar.info("Click 'Create GPT Tasks' to proceed")
            if st.sidebar.button("Create GPT Tasks"):
                for key, checked in checkboxes.items():
                    if checked:
                        st.session_state.prompts.append({
                            "prompt": pre_canned_prompts[summary_type.lower().replace(" ", "_")][key]["prompt"],
                            "model": "gpt-4o-mini",
                            "heading": pre_canned_prompts[summary_type.lower().replace(" ", "_")][key]["heading"]
                        })

        for i, prompt_info in enumerate(st.session_state.prompts):
            with st.expander(f"GPT Task {i+1} - {prompt_info['heading']}", expanded=True):
                st.info("Update the pre-canned prompt to customize!")
                prompt_info["model"] = st.text_input("Model", value=prompt_info["model"], key=f"model_{i}")
                prompt_info["prompt"] = st.text_area("Prompt", value=prompt_info["prompt"], key=f"prompt_{i}")
                if st.button("Remove GPT Task", key=f"remove_gpt_task_{i}"):
                    st.session_state.prompts.pop(i)
                    break

        if st.session_state.prompts:
            st.info("Click generate to create your document!")
            st.markdown(
                """
                <style>
                .blue-button button {
                    background-color: #007BFF !important;
                    color: white !important;
                }
                </style>
                """, unsafe_allow_html=True
            )
            if st.button("Generate", key="generate"):
                minutes = {prompt_info["heading"]: generate_response(st.session_state.transcription, prompt_info["model"], prompt_info["prompt"]) for prompt_info in st.session_state.prompts}
                st.session_state.generated_minutes = minutes

        if 'generated_minutes' in st.session_state:
            with st.expander("Generated Minutes", expanded=True):
                for key, value in st.session_state.generated_minutes.items():
                    st.write(f"**{key}**")
                    st.write(value)

                docx_file = save_as_docx(st.session_state.generated_minutes)

                st.info("Click download to get a docx file of your document!")
                st.download_button(
                    label="Download Meeting Minutes",
                    data=docx_file,
                    file_name="meeting_minutes.docx",
                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                )

            if "Action Items" in st.session_state.generated_minutes:
                with st.expander("Action Items", expanded=True):
                    st.subheader("Action Items")
                    action_items = st.session_state.generated_minutes["Action Items"]
                    st.info("Check boxes to generate documents from tasks!")
                    action_items_list = [item for item in action_items.split('\n') if item]

                    action_items_dict = {}
                    parent_task = None

                    for item in action_items_list:
                        if item.startswith("    "):
                            if parent_task:
                                action_items_dict[parent_task].append(item.strip())
                        else:
                            parent_task = item.strip()
                            action_items_dict[parent_task] = []

                    grid_df = pd.DataFrame({
                        "Task Number": range(1, len(action_items_dict) + 1),
                        "Task": list(action_items_dict.keys()),
                        "Draft Email": False,
                        "Draft Slack": False,
                        "Draft Memo": False
                    })

                    gb = GridOptionsBuilder.from_dataframe(grid_df)
                    gb.configure_column("Draft Email", editable=True, cellEditor="agCheckboxCellEditor")
                    gb.configure_column("Draft Slack", editable=True, cellEditor="agCheckboxCellEditor")
                    gb.configure_column("Draft Memo", editable=True, cellEditor="agCheckboxCellEditor")
                    gb.configure_pagination()
                    gb.configure_default_column(editable=True, resizable=True)
                    grid_options = gb.build()

                    grid_response = AgGrid(grid_df, gridOptions=grid_options, height=300, fit_columns_on_grid_load=True, update_mode=GridUpdateMode.MODEL_CHANGED)

                    if isinstance(grid_response['data'], pd.DataFrame):
                        for _, row in grid_response['data'].iterrows():
                            task_num = row['Task Number']
                            if row["Draft Email"]:
                                st.session_state[f"email_prompt_{task_num}"] = f"Draft an email for the following action item: {row['Task']}"
                            if row["Draft Slack"]:
                                st.session_state[f"slack_prompt_{task_num}"] = f"Draft a Slack message for the following action item: {row['Task']}"
                            if row["Draft Memo"]:
                                st.session_state[f"memo_prompt_{task_num}"] = f"Draft a memo for the following action item: {row['Task']}"

                    for key, value in st.session_state.items():
                        if key.startswith("email_prompt_") or key.startswith("slack_prompt_") or key.startswith("memo_prompt_"):
                            task_num = key.split('_')[-1]
                            st.subheader(f"{key.split('_')[0].capitalize()} Draft for Task {task_num}")
                            st.write(value)
                            if st.button(f"Generate {key.split('_')[0].capitalize()} for Task {task_num}"):
                                draft = generate_response(st.session_state.transcription, "gpt-4o-mini", value)
                                st.write(draft)

if __name__ == "__main__":
    main()
